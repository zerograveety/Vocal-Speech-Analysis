"""
PPT analysis utilities:
 - Parse PPTX slides (text and images)
 - Compute text density and visual balance
 - Estimate grammar acceptability with a Transformer classifier (CoLA)
"""

from __future__ import annotations

from typing import Dict, List, Tuple
from dataclasses import dataclass

from pptx import Presentation
import numpy as np

# Optional heavy deps. We fallback gracefully if unavailable.
try:
    import torch  # type: ignore
    from transformers import AutoTokenizer, AutoModelForSequenceClassification  # type: ignore
    _NLP_DEPS_AVAILABLE = True
except Exception:
    torch = None  # type: ignore
    AutoTokenizer = None  # type: ignore
    AutoModelForSequenceClassification = None  # type: ignore
    _NLP_DEPS_AVAILABLE = False


MODEL_NAME = "textattack/bert-base-uncased-CoLA"


@dataclass
class SlideStats:
    num_words: int
    num_images: int


def _extract_slide_text_and_images(prs: Presentation) -> List[SlideStats]:
    stats: List[SlideStats] = []
    for slide in prs.slides:
        text_runs: List[str] = []
        image_count = 0
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                image_count += 1
            if shape.shape_type == 1 and hasattr(shape, "text_frame"):
                # AutoShape with text
                try:
                    for p in shape.text_frame.paragraphs:
                        text_runs.append(" ".join(run.text for run in p.runs))
                except Exception:
                    pass
        num_words = sum(len(t.split()) for t in text_runs)
        stats.append(SlideStats(num_words=num_words, num_images=image_count))
    return stats


def _load_cola_model():
    if not _NLP_DEPS_AVAILABLE:
        return None, None
    tokenizer = AutoTokenizer.from_pretrained(MODEL_NAME)
    model = AutoModelForSequenceClassification.from_pretrained(MODEL_NAME)
    model.eval()
    return tokenizer, model


def _grammar_score(text: str, tokenizer, model) -> float:
    if not text.strip():
        return 0.5
    if tokenizer is None or model is None or torch is None:
        # Lightweight heuristic fallback: punctuation density proxy (0-1)
        total = max(len(text), 1)
        punct = sum(1 for c in text if c in ".!?;:")
        ratio = min(punct / total * 10.0, 1.0)
        return float(round(0.3 + 0.4 * ratio, 3))
    inputs = tokenizer(text, truncation=True, max_length=256, return_tensors="pt")
    with torch.no_grad():
        outputs = model(**inputs)
        probs = torch.softmax(outputs.logits, dim=-1)
        acceptability = probs[0, 1].item()
    return float(acceptability)


def analyze_pptx(pptx_path: str) -> Dict[str, object]:
    prs = Presentation(pptx_path)
    stats = _extract_slide_text_and_images(prs)

    total_words = sum(s.num_words for s in stats)
    total_images = sum(s.num_images for s in stats)
    num_slides = len(stats) if stats else 1

    avg_words_per_slide = total_words / num_slides
    avg_images_per_slide = total_images / num_slides

    # Text density heuristic
    if avg_words_per_slide > 80:
        slide_quality = "Too much text"
    elif avg_words_per_slide < 20:
        slide_quality = "Too little text"
    else:
        slide_quality = "Balanced"

    # Visual balance heuristic
    if avg_images_per_slide >= 2 and avg_words_per_slide <= 60:
        visual_balance = "Visual-heavy and balanced"
    elif avg_images_per_slide == 0 and avg_words_per_slide > 80:
        visual_balance = "Text-heavy"
    else:
        visual_balance = "Moderate"

    # Grammar score over concatenated slide text (first N tokens)
    all_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                all_text.append(shape.text)
    combined_text = "\n".join(all_text)

    tokenizer, model = _load_cola_model()
    grammar = _grammar_score(combined_text, tokenizer, model)

    return {
        "slide_quality": slide_quality,
        "grammar_score": round(grammar, 3),
        "visual_balance": visual_balance,
        "details": {
            "avg_words_per_slide": round(avg_words_per_slide, 1),
            "avg_images_per_slide": round(avg_images_per_slide, 2),
            "num_slides": num_slides,
        },
    }


